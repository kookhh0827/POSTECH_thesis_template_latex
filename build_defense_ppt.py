"""
Master's thesis defense slides — Option A (18 slides, ~20 min)
Hyunho Kook -- Stabilizing Direct Training of SNNs

Design follows reference deck MP-Init_TrSG_Defense.pptx:
  - Navy / teal / amber / coral palette
  - Serif title (Cambria), sans body (Calibri)
  - Decorative circles on cover/conclusion
  - Section label · slide title · content
  - Audience: advisor + ML/CS people who don't know SNN
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
import os

ASSETS = "ppt_assets"
GEN    = "ppt_assets/gen"

# 16:9 widescreen
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Palette (from reference) ──────────────────────────────────────
NAVY     = RGBColor(0x0F, 0x25, 0x40)
NAVY_DK  = RGBColor(0x09, 0x1A, 0x30)
TEAL     = RGBColor(0x1D, 0x7E, 0x7E)
TEAL_LT  = RGBColor(0xE5, 0xF5, 0xF0)
AMBER    = RGBColor(0xE0, 0xB8, 0x5F)
AMBER_LT = RGBColor(0xFA, 0xEE, 0xCB)
CORAL    = RGBColor(0xE5, 0x69, 0x45)
CORAL_LT = RGBColor(0xFA, 0xEE, 0xE0)
INK      = RGBColor(0x1A, 0x1A, 0x1A)
GRAY     = RGBColor(0x7A, 0x7A, 0x7A)
GRAY_LT  = RGBColor(0xE5, 0xE5, 0xE0)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CREAM    = RGBColor(0xF8, 0xF7, 0xF2)

SERIF = "Cambria"
SANS  = "Calibri"

# ─── Helpers ───────────────────────────────────────────────────────
def add_blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element); spTree.insert(2, bg._element)
    return bg

def add_circle(slide, x, y, d, color):
    """Decorative circle for title/cover slides."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = color
    c.line.fill.background()
    return c

def text(slide, x, y, w, h, txt, *, font=SANS, size=18, bold=False, italic=False,
         color=INK, align=None, anchor=None, spacing=None):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = txt
    p.font.name = font
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align: p.alignment = align
    if anchor: tf.vertical_anchor = anchor
    if spacing:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            rPr.set('spc', str(spacing))
    return tx

def section_label(slide, txt):
    """Top section label, ALL CAPS letter-spaced teal."""
    text(slide, 0.5, 0.4, 12, 0.3, txt.upper(),
         font=SANS, size=11, bold=True, color=TEAL, spacing=300)

def slide_title(slide, txt, sub=None):
    text(slide, 0.5, 0.85, 12.3, 0.85, txt,
         font=SERIF, size=30, bold=True, color=NAVY)
    if sub:
        text(slide, 0.5, 1.75, 12.3, 0.4, sub,
             font=SANS, size=14, italic=True, color=GRAY)

def page_footer(slide, n, total=20):
    text(slide, 0.5, 7.10, 6, 0.3, "Defense — Hyunho Kook",
         font=SANS, size=10, color=GRAY)
    text(slide, 7.0, 7.10, 5.8, 0.3, f"{n} / {total}",
         font=SANS, size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def amber_underline(slide, x, y, w=0.7):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.04))
    bar.fill.solid(); bar.fill.fore_color.rgb = AMBER
    bar.line.fill.background()
    return bar

def callout_box(slide, x, y, w, h, *, fill=NAVY, accent=AMBER,
                title=None, body=None, title_size=14, body_size=12,
                title_color=AMBER, body_color=WHITE):
    # vertical accent line
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(0.06), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # background box
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x + 0.06), Inches(y), Inches(w - 0.06), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.fill.background()
    # text
    if title or body:
        tx = slide.shapes.add_textbox(
            Inches(x + 0.35), Inches(y + 0.05),
            Inches(w - 0.5), Inches(h - 0.1))
        tf = tx.text_frame; tf.word_wrap = True
        # Vertically center the text within the callout box
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        if title:
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = SANS; p.font.size = Pt(title_size); p.font.bold = True
            p.font.color.rgb = title_color
        if body:
            p2 = tf.add_paragraph() if title else tf.paragraphs[0]
            p2.text = body
            p2.font.name = SANS; p2.font.size = Pt(body_size)
            p2.font.color.rgb = body_color
            p2.space_before = Pt(6) if title else Pt(0)

def card(slide, x, y, w, h, *, fill=WHITE, border=TEAL, border_w=1.5):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = border; rect.line.width = Pt(border_w)
    return rect

def bullets(slide, x, y, w, h, items, *, size=14, color=INK,
            bullet="•", spacing=6):
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.name = SANS; p.font.size = Pt(size); p.font.color.rgb = color
        if i > 0: p.space_before = Pt(spacing)


def img(slide, path, x, y, w=None, h=None):
    if not os.path.exists(path):
        return None
    kw = {}
    if w is not None: kw["width"] = Inches(w)
    if h is not None: kw["height"] = Inches(h)
    return slide.shapes.add_picture(path, Inches(x), Inches(y), **kw)


# ════════════════════════════════════════════════════════════════════
# S1 — Title
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, NAVY)
# Decorative circles
add_circle(s, -2.0, -2.0, 5.5, NAVY_DK)
add_circle(s, 10.5, 4.5, 4.0, NAVY_DK)

text(s, 0.7, 0.9, 10, 0.4, "MASTER'S THESIS DEFENSE",
     font=SANS, size=12, bold=True, color=AMBER, spacing=400)
text(s, 0.7, 1.7, 12, 1.3,
     "Enhancing Stability in Direct Training",
     font=SERIF, size=44, bold=True, color=WHITE)
text(s, 0.7, 2.55, 12, 1.3,
     "of Spiking Neural Networks",
     font=SERIF, size=44, bold=True, color=WHITE)
text(s, 0.7, 3.7, 12, 0.5,
     "Membrane Potential Initialization & Threshold-Robust Surrogate Gradients",
     font=SERIF, size=16, italic=True, color=AMBER)

# Amber underline
amber_underline(s, 0.7, 4.45, w=1.0)

text(s, 0.7, 4.7, 9, 0.4, "Hyunho Kook  (국 현 호)",
     font=SANS, size=18, bold=True, color=WHITE)
text(s, 0.7, 5.2, 9, 0.35, "Advisor: Prof. Eunhyeok Park",
     font=SANS, size=12, color=RGBColor(0xCB, 0xD8, 0xE0))
text(s, 0.7, 5.55, 9, 0.35, "Committee: Prof. Jungseul Ok · Prof. Jae-Joon Kim",
     font=SANS, size=12, color=RGBColor(0xCB, 0xD8, 0xE0))

text(s, 0.7, 7.05, 12, 0.3,
     "Department of Computer Science and Engineering · POSTECH · January 2026",
     font=SANS, size=10, color=RGBColor(0x80, 0x95, 0xA8))


# ════════════════════════════════════════════════════════════════════
# S2 — Why SNNs?
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Motivation")
slide_title(s, "Why Spiking Neural Networks?")

# Left column: 3 rows of bullet points
y = 2.4
text(s, 0.5, y, 7.2, 0.45,
     "DNNs are powerful but expensive",
     font=SANS, size=15, bold=True, color=NAVY)
text(s, 0.5, y+0.5, 7.2, 1.0,
     "Training and serving large models routinely consumes megawatt-scale power; "
     "even edge inference is hard under tight energy budgets.",
     font=SANS, size=12, color=INK)

text(s, 0.5, y+1.6, 7.2, 0.45,
     "SNNs: an event-driven alternative",
     font=SANS, size=15, bold=True, color=NAVY)
text(s, 0.5, y+2.1, 7.2, 1.0,
     "Information flows as discrete spikes over time. Paired with neuromorphic hardware, "
     "computation skips inactive neurons — orders of magnitude lower energy.",
     font=SANS, size=12, color=INK)

text(s, 0.5, y+3.2, 7.2, 0.45,
     "But direct training is unstable",
     font=SANS, size=15, bold=True, color=CORAL)
text(s, 0.5, y+3.7, 7.2, 1.0,
     "Recent direct-training methods narrowed the accuracy gap, yet two persistent instabilities "
     "still hold SNNs back.",
     font=SANS, size=12, color=INK)

# Right column: dark callout panel
panel_x, panel_y, panel_w, panel_h = 8.5, 2.3, 4.4, 4.4
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(panel_x), Inches(panel_y), Inches(panel_w), Inches(panel_h))
box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.fill.background()

text(s, panel_x + 0.3, panel_y + 0.25, panel_w - 0.5, 0.3,
     "THIS THESIS", font=SANS, size=10, bold=True, color=AMBER, spacing=400)
text(s, panel_x + 0.3, panel_y + 0.7, panel_w - 0.5, 0.5,
     "Two instabilities,",
     font=SERIF, size=22, bold=True, color=WHITE)
text(s, panel_x + 0.3, panel_y + 1.15, panel_w - 0.5, 0.5,
     "two principled fixes.",
     font=SERIF, size=22, bold=True, color=WHITE)

# Two numbered items inside panel
nx, ny1, ny2 = panel_x + 0.3, panel_y + 2.0, panel_y + 3.2
# accent bar 1
ab1 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(nx), Inches(ny1), Inches(0.04), Inches(0.85))
ab1.fill.solid(); ab1.fill.fore_color.rgb = TEAL; ab1.line.fill.background()
text(s, nx + 0.15, ny1, panel_w - 0.6, 0.3,
     "①  Temporal Covariate Shift",
     font=SANS, size=13, bold=True, color=WHITE)
text(s, nx + 0.15, ny1 + 0.35, panel_w - 0.6, 0.5,
     "in the membrane potential — invisible to existing remedies.",
     font=SANS, size=10, color=RGBColor(0xCB, 0xD8, 0xE0))
# accent bar 2
ab2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(nx), Inches(ny2), Inches(0.04), Inches(0.95))
ab2.fill.solid(); ab2.fill.fore_color.rgb = AMBER; ab2.line.fill.background()
text(s, nx + 0.15, ny2, panel_w - 0.6, 0.3,
     "②  Gradient pathology",
     font=SANS, size=13, bold=True, color=WHITE)
text(s, nx + 0.15, ny2 + 0.35, panel_w - 0.6, 0.6,
     "when the firing threshold becomes a learnable parameter — never explicitly characterized before.",
     font=SANS, size=10, color=RGBColor(0xCB, 0xD8, 0xE0))

page_footer(s, 2)


# ════════════════════════════════════════════════════════════════════
# S3 — SNNs in 60 seconds (LIF + surrogate gradient)
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Background · 1/2")
slide_title(s, "Spiking Neurons in 60 Seconds",
            "Three discrete-time updates, plus one trick to get gradients through them")

# Left: LIF equations
text(s, 0.5, 2.5, 6.2, 0.4,
     "Leaky Integrate-and-Fire neuron",
     font=SANS, size=14, bold=True, color=NAVY)

eq_y = 3.0
def equation_box(slide, x, y, w, h, label, eq):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CREAM
    box.line.color.rgb = GRAY_LT; box.line.width = Pt(0.5)
    text(slide, x + 0.15, y + 0.05, 1.4, 0.3,
         label, font=SANS, size=10, bold=True, color=TEAL)
    text(slide, x + 1.5, y + 0.05, w - 1.6, h - 0.1, eq,
         font="Cambria Math", size=12, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)

equation_box(s, 0.5, eq_y,        6.2, 0.55,
             "Integrate", "M[t] = (1−1/τ)·U[t−1]  +  (1/τ)·I_in[t]")
equation_box(s, 0.5, eq_y + 0.65, 6.2, 0.55,
             "Fire",      "S[t] = H( M[t] − V_thr ) ∈ {0, 1}")
equation_box(s, 0.5, eq_y + 1.30, 6.2, 0.55,
             "Reset",     "U[t] = M[t] − V_thr · S[t]   (soft reset)")

text(s, 0.5, eq_y + 1.95, 6.2, 0.4,
     "Two learnable parameters:  τ (leakage)  ·  V_thr (threshold).",
     font=SANS, size=11, italic=True, color=GRAY)

# Right: LIF dynamics chart + surrogate gradient
img(s, GEN + "/lif_dynamics.png", 7.0, 2.5, w=5.9)
text(s, 7.0, 5.05, 5.9, 0.3,
     "M[t] integrates inputs, fires at V_thr, resets — repeat.",
     font=SANS, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

text(s, 7.0, 5.45, 5.9, 0.4,
     "Heaviside H is not differentiable",
     font=SANS, size=12, bold=True, color=CORAL)
text(s, 7.0, 5.85, 5.9, 0.45,
     "Surrogate gradient: backward pass uses smooth f'(x) ≈ H'(x).",
     font=SANS, size=11, color=INK)

# Bottom callout (single-line title only — keep tight)
callout_box(s, 0.5, 6.2, 12.3, 0.45,
            title="Everything in this thesis turns on M[t] and V_thr  ·  TCS at M[t]  ·  pathology at V_thr",
            title_size=12)

page_footer(s, 3)


# ════════════════════════════════════════════════════════════════════
# S4 — Diagnosis ① TCS in membrane potential
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Background · 2/2")
slide_title(s, "Diagnosis ①: TCS Lives in the Membrane Potential")

# Left: prose
text(s, 0.5, 2.5, 6.5, 0.4,
     "What everyone notices",
     font=SANS, size=14, bold=True, color=NAVY)
text(s, 0.5, 2.95, 6.5, 1.5,
     "Activation distributions drift across timesteps. Batch-norm statistics estimated\n"
     "at one timestep no longer fit at another — accuracy drops.",
     font=SANS, size=12, color=INK)

text(s, 0.5, 4.2, 6.5, 0.4,
     "What existing fixes do",
     font=SANS, size=14, bold=True, color=NAVY)
text(s, 0.5, 4.65, 6.5, 1.0,
     "TEBN, TAB add per-timestep BN parameters.\n"
     "More compute, more parameters — and they don't reach the cause.",
     font=SANS, size=12, color=INK)

# Bottom callout — the claim
callout_box(s, 0.5, 6.0, 6.5, 0.95,
            fill=NAVY, accent=CORAL,
            title="Our claim",
            body="TCS originates one level deeper — in M[t] itself.",
            title_color=CORAL, title_size=13, body_size=12)

# Right: figure
img(s, ASSETS + "/density_plot_membrane_potential_tdbn.png", 7.5, 2.4, w=5.4)
text(s, 7.5, 6.3, 5.4, 0.3,
     "M[t] under tdBN — distribution clearly drifts as t grows",
     font=SANS, size=10, italic=True, color=CORAL, align=PP_ALIGN.CENTER)

page_footer(s, 4)


# ════════════════════════════════════════════════════════════════════
# S5 — Diagnosis ② V_thr training breaks gradient
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Setup")
slide_title(s, "Diagnosis ②: Trainable V_thr Breaks the Gradient",
            "PLIF, LTMD, DIET-SNN all let V_thr learn — but the surrogate gradient is scale-sensitive")

# Two failure scenarios
fy = 2.5

# Left scenario
left_x, right_x, w, h = 0.5, 6.95, 5.85, 2.6
rect_l = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(left_x), Inches(fy), Inches(w), Inches(h))
rect_l.fill.solid(); rect_l.fill.fore_color.rgb = CORAL_LT
rect_l.line.color.rgb = CORAL; rect_l.line.width = Pt(1.2)
text(s, left_x + 0.3, fy + 0.2, w - 0.5, 0.4,
     "V_thr ≪ 1   (small threshold)",
     font=SANS, size=14, bold=True, color=CORAL)
bullets(s, left_x + 0.3, fy + 0.7, w - 0.5, 1.8, [
    "AS-SG floods — fixed window covers almost the entire distribution",
    "RS-SG explodes — the 1/V_thr factor blows the magnitude up",
    "Real example: at V_thr = 0.1, RS-SG diverges to NaN within 100 epochs",
], size=11, color=INK, bullet="▸")

# Right scenario
rect_r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(right_x), Inches(fy), Inches(w), Inches(h))
rect_r.fill.solid(); rect_r.fill.fore_color.rgb = CORAL_LT
rect_r.line.color.rgb = CORAL; rect_r.line.width = Pt(1.2)
text(s, right_x + 0.3, fy + 0.2, w - 0.5, 0.4,
     "V_thr ≫ 1   (large threshold)",
     font=SANS, size=14, bold=True, color=CORAL)
bullets(s, right_x + 0.3, fy + 0.7, w - 0.5, 1.8, [
    "AS-SG starves — fixed window covers a tiny tail of the distribution",
    "RS-SG vanishes — the same 1/V_thr factor shrinks the magnitude",
    "Real example: at V_thr = 2.0, AS-SG hits 0% active neurons by epoch 100",
], size=11, color=INK, bullet="▸")

# Bottom: real-world emphasis
callout_box(s, 0.5, 5.4, 12.3, 1.4,
            fill=NAVY, accent=AMBER,
            title="These are not corner cases.",
            body="Trainable V_thr in practice drifts well below 1 in early layers and above 2 in deep ones —\n"
                 "exactly the regimes where AS-SG and RS-SG both fail. We'll quantify in Part II.",
            title_color=AMBER, body_color=WHITE,
            title_size=14, body_size=11)

page_footer(s, 5)


# ════════════════════════════════════════════════════════════════════
# S6 — Theory: Markov chain
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part I · MP-Init")
slide_title(s, "Why Does the Membrane Potential Drift?",
            "Because it's a Markov chain pulled toward a unique distribution")

# Left: theorem + recurrence (more compact)
callout_box(s, 0.5, 2.5, 7.2, 0.7,
            fill=CREAM, accent=TEAL,
            title="The recurrence is a Markov chain",
            body="U[t+1] = f( U[t], I_in[t+1] ) — Markov by construction;  i.i.d. inputs ⇒ time-homogeneous",
            title_color=TEAL, body_color=INK,
            title_size=12, body_size=11)

callout_box(s, 0.5, 3.4, 7.2, 1.6,
            fill=NAVY, accent=AMBER,
            title="Theorem  (Doeblin minorization)",
            body="Under (i.i.d. inputs) + (bounded U[t]),\n"
                 "the chain has a unique stationary π.  Moreover,\n\n"
                 "‖ Pⁿ(x, ·) − π(·) ‖_TV   ≤   C · (1 − ε)ⁿ",
            title_color=AMBER, body_color=WHITE,
            title_size=12, body_size=11)

# Implications
text(s, 0.5, 5.2, 7.2, 0.4,
     "Implications",
     font=SANS, size=13, bold=True, color=CORAL)
bullets(s, 0.5, 5.6, 7.2, 1.4, [
    "Initial state ≠ π  ⇒  drift is inevitable",
    "Standard U[0] = 0 is far from π",
    "Don't fix BN — fix the initial state",
], size=11, color=INK, bullet="▸")

# Right: convergence visualization
img(s, GEN + "/convergence.png", 8.0, 2.5, w=5.0)
text(s, 8.0, 5.4, 5.0, 0.4,
     "Distribution of U[t] for t = 0, 1, 2, 3, 4, 5, 8, 12",
     font=SANS, size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
text(s, 8.0, 5.8, 5.0, 0.7,
     "Starting from U[0]=0 (sharp red), the distribution exponentially "
     "approaches the stationary π (dashed navy).",
     font=SANS, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

page_footer(s, 6)


# ════════════════════════════════════════════════════════════════════
# S7 — MP-Init: Algorithm
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part I · MP-Init")
slide_title(s, "MP-Init — Initialize at the Stationary Mean",
            "U[T] is the closest sample to π → use a running mean of E[U[T]]")

# Left: lemma + estimate
text(s, 0.5, 2.5, 6.5, 0.4,
     "Lemma — best constant is the mean",
     font=SANS, size=14, bold=True, color=NAVY)
callout_box(s, 0.5, 2.95, 6.5, 0.7,
            fill=CREAM, accent=TEAL,
            body="argmin_c   E[ (π − c)² ]   =   E[π]",
            body_color=INK, body_size=14, title_size=12)

text(s, 0.5, 3.95, 6.5, 0.4,
     "Practical estimator",
     font=SANS, size=14, bold=True, color=NAVY)
text(s, 0.5, 4.4, 6.5, 1.5,
     "U[T] is the closest sample to π. Maintain a per-layer running mean\n"
     "μ ← β · μ  +  (1 − β) · mean( U[T] )    over training, β = 0.9.\n"
     "At every batch, set U[0] = μ.",
     font=SANS, size=12, color=INK)

# Right: pseudocode
code_x = 7.4
text(s, code_x, 2.5, 5.4, 0.4,
     "Algorithm  (training-only)",
     font=SANS, size=14, bold=True, color=NAVY)

code_box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(code_x), Inches(2.95), Inches(5.4), Inches(2.7))
code_box.fill.solid(); code_box.fill.fore_color.rgb = NAVY
code_box.line.fill.background()

code_lines = [
    "init  μ_l ← 0   for each layer l",
    "",
    "for each batch:",
    "    U_l[0] ← μ_l",
    "    run T LIF steps as usual",
    "    m ← mean( U_l[T] of active neurons )",
    "    μ_l ← β · μ_l  +  (1 − β) · m",
]
ctx = s.shapes.add_textbox(Inches(code_x + 0.25), Inches(3.05),
                            Inches(5.0), Inches(2.5))
tf = ctx.text_frame; tf.word_wrap = True
for i, line in enumerate(code_lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.name = "Consolas"; p.font.size = Pt(11.5); p.font.color.rgb = WHITE

# Bottom strip
callout_box(s, 0.5, 6.05, 12.3, 0.9,
            fill=TEAL, accent=AMBER,
            title="Inference: μ frozen → standard LIF dynamics",
            body="Zero changes to BN, no per-timestep parameters, full neuromorphic-hardware compatibility.",
            title_color=WHITE, body_color=RGBColor(0xE0, 0xF0, 0xEA),
            title_size=13, body_size=11)

page_footer(s, 7)


# ════════════════════════════════════════════════════════════════════
# S8 — MP-Init: Result + Mechanism
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part I · MP-Init")
slide_title(s, "MP-Init Works — at Both the Distribution and the Metric")

# Top row: 4 panels showing distribution alignment
text(s, 0.5, 2.4, 12.3, 0.4,
     "Membrane potential distributions across timesteps  (ResNet-19 / CIFAR-100)",
     font=SANS, size=11, italic=True, color=GRAY)
img(s, ASSETS + "/density_plot_membrane_potential_tdbn.png", 0.5, 2.8, w=2.95)
img(s, ASSETS + "/density_plot_membrane_potential_tebn.png", 3.6, 2.8, w=2.95)
img(s, ASSETS + "/density_plot_membrane_potential_tab.png",  6.7, 2.8, w=2.95)
img(s, ASSETS + "/density_plot_membrane_potential_mpinit.png", 9.8, 2.8, w=2.95)
text(s, 0.5, 4.65, 2.95, 0.3, "tdBN", font=SANS, size=10, color=GRAY, align=PP_ALIGN.CENTER)
text(s, 3.6, 4.65, 2.95, 0.3, "TEBN", font=SANS, size=10, color=GRAY, align=PP_ALIGN.CENTER)
text(s, 6.7, 4.65, 2.95, 0.3, "TAB",  font=SANS, size=10, color=GRAY, align=PP_ALIGN.CENTER)
text(s, 9.8, 4.65, 2.95, 0.3, "MP-Init (ours)",
     font=SANS, size=10, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

# Bottom: gain table card
text(s, 0.5, 4.95, 12.3, 0.4,
     "Consistent accuracy gain on every baseline",
     font=SANS, size=14, bold=True, color=NAVY)

# Gain table — native PPT table for cleaner grid + edit-ability
gx, gy = 0.5, 5.30
gain_rows = [
    ("CIFAR-100  ·  T=2  ·  tdBN",      "72.35", "74.01"),
    ("CIFAR-100  ·  T=4  ·  tdBN",      "75.55", "76.09"),
    ("CIFAR-100  ·  T=4  ·  TEBN",      "75.96", "76.45"),
    ("CIFAR-100  ·  T=4  ·  TAB",       "76.25", "77.24"),
    ("DVS-CIFAR10  ·  T=10  ·  tdBN",   "76.60", "77.37"),
]
gain_tbl_shape = s.shapes.add_table(
    len(gain_rows) + 1, 3,
    Inches(gx), Inches(gy),
    Inches(7.2), Inches(1.55))
gain_tbl = gain_tbl_shape.table
gain_tbl.columns[0].width = Inches(4.2)
gain_tbl.columns[1].width = Inches(1.5)
gain_tbl.columns[2].width = Inches(1.5)
gain_tbl.rows[0].height = Inches(0.32)
for r in range(1, len(gain_rows) + 1):
    gain_tbl.rows[r].height = Inches(0.245)

# Header row
for j, (h, alignj) in enumerate([("Setting", PP_ALIGN.LEFT),
                                  ("Baseline", PP_ALIGN.CENTER),
                                  ("+ MP-Init", PP_ALIGN.CENTER)]):
    cell = gain_tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_LT
    p = cell.text_frame.paragraphs[0]
    p.font.name = SANS; p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = alignj
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.15)
    cell.margin_right = Inches(0.15)

# Data rows with zebra striping
for i, (label, base, ours) in enumerate(gain_rows, start=1):
    bg = CREAM if i % 2 == 0 else WHITE
    for j, (val, font_color, is_bold, alignj) in enumerate([
        (label, INK,  False, PP_ALIGN.LEFT),
        (base,  GRAY, False, PP_ALIGN.CENTER),
        (ours,  TEAL, True,  PP_ALIGN.CENTER),
    ]):
        cell = gain_tbl.cell(i, j)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.font.name = SANS; p.font.size = Pt(11); p.font.bold = is_bold
        p.font.color.rgb = font_color
        p.alignment = alignj
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.15)
        cell.margin_right = Inches(0.15)

# Right: highlighted +1.66
hx = 8.0
hbox = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(hx), Inches(5.30), Inches(4.85), Inches(1.55))
hbox.fill.solid(); hbox.fill.fore_color.rgb = AMBER_LT
hbox.line.color.rgb = AMBER; hbox.line.width = Pt(1.2)
text(s, hx + 0.25, 5.45, 4.5, 0.55,
     "+1.66 %pt at T = 2",
     font=SERIF, size=24, bold=True, color=NAVY)
text(s, hx + 0.25, 6.05, 4.5, 0.7,
     "Biggest gain in the most latency-critical regime —\nexactly where SNN energy efficiency matters.",
     font=SANS, size=10, color=INK)

page_footer(s, 8)


# ════════════════════════════════════════════════════════════════════
# S9 — Why MP-Init Works  (mechanism)
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part I · MP-Init")
slide_title(s, "Why MP-Init Works — A Closer Look",
            "Aligning distributions aligns gradients — and improves the firing-rate frontier")

# Left column: gradient cosine similarity
text(s, 0.5, 2.4, 6.0, 0.4,
     "Gradient cosine similarity across timesteps",
     font=SANS, size=13, bold=True, color=NAVY)
# grad_collision aspect ≈ 2.01.  h=2.0 → w=4.02. Center in left half.
img(s, ASSETS + "/grad_collision.png", 1.49, 2.85, h=2.0)
text(s, 0.5, 4.95, 6.0, 0.3,
     "ResNet-19 · CIFAR-100 · 4 timesteps",
     font=SANS, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
bullets(s, 0.5, 5.3, 6.0, 1.4, [
    "Without MP-Init: gradients at t=1 conflict with later steps (cos 0.07–0.18)",
    "With MP-Init: temporal ensemble realigned (cos ≈ 0.5)",
    "Stable optimization — not just stable inference",
], size=11, color=INK, bullet="▸")

# Right column: Pareto frontier
text(s, 6.85, 2.4, 6.0, 0.4,
     "Accuracy vs. firing rate  (Pareto frontier)",
     font=SANS, size=13, bold=True, color=NAVY)
# firing_rate_vs_accuracy aspect ≈ 2.85.  h=2.0 → w=5.70.
img(s, ASSETS + "/firing_rate_vs_accuracy.png", 7.0, 2.85, h=2.0)
text(s, 6.85, 4.95, 6.0, 0.3,
     "ResNet-19 · CIFAR-100 with sparsity regularizer",
     font=SANS, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
bullets(s, 6.85, 5.3, 6.0, 1.4, [
    "MP-Init dominates baseline at every firing rate",
    "Even at 6% firing rate, beats all baseline settings",
    "Gains come from spiking BETTER, not spiking MORE",
], size=11, color=INK, bullet="▸")

# Bottom takeaway
callout_box(s, 0.5, 6.65, 12.3, 0.45,
            fill=TEAL, accent=AMBER,
            title="Higher accuracy from stabler training — efficiency-preserving improvements, not accuracy-via-spend.",
            title_color=WHITE, title_size=12)

page_footer(s, 9)


# ════════════════════════════════════════════════════════════════════
# S10 — Two SG families
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part II · TrSG")
slide_title(s, "Two Surrogate-Gradient Families",
            "AS-SG and RS-SG — look similar, behave very differently when V_thr is trainable")

# Recap of surrogate gradient (compact)
callout_box(s, 0.5, 2.4, 12.3, 0.55,
            fill=CREAM, accent=AMBER,
            title="Recap — Surrogate Gradient: replace H'(x) with smooth f'(x).  Choice of x differs below.",
            title_color=NAVY, title_size=12)

# Left side: definitions stacked vertically
card(s, 0.5, 3.15, 5.2, 1.65, fill=CREAM, border=NAVY, border_w=1.5)
text(s, 0.7, 3.25, 4.8, 0.35, "AS-SG  ·  Absolute-Scale",
     font=SANS, size=13, bold=True, color=NAVY)
text(s, 0.7, 3.65, 4.8, 0.5,
     "x = M[t] − V_thr",
     font="Cambria Math", size=18, color=NAVY)
text(s, 0.7, 4.30, 4.8, 0.4,
     "Window of FIXED width γ, centered at V_thr",
     font=SANS, size=11, italic=True, color=CORAL)

card(s, 0.5, 4.95, 5.2, 1.65, fill=CREAM, border=NAVY, border_w=1.5)
text(s, 0.7, 5.05, 4.8, 0.35, "RS-SG  ·  Relative-Scale",
     font=SANS, size=13, bold=True, color=NAVY)
text(s, 0.7, 5.45, 4.8, 0.5,
     "x = M[t] / V_thr − 1",
     font="Cambria Math", size=18, color=NAVY)
text(s, 0.7, 6.10, 4.8, 0.4,
     "Window SCALES with V_thr — but magnitude / V_thr",
     font=SANS, size=11, italic=True, color=AMBER)

# Right side: BIG window visualization (now properly sized)
text(s, 5.95, 3.15, 7.0, 0.3,
     "Gradient-window vs membrane distribution",
     font=SANS, size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)
# sg_windows is 8:3 aspect.  At w=7.0 → h=2.62.  Position y=3.5 → ends 6.12.
img(s, GEN + "/sg_windows.png", 5.95, 3.5, w=7.0)

# Bottom strip (sits above the page footer)
callout_box(s, 0.5, 6.55, 12.3, 0.4,
            fill=NAVY, accent=CORAL,
            title="Same when V_thr is fixed. Once V_thr trains, AS-SG and RS-SG break in opposite ways.",
            title_color=AMBER, title_size=11)

page_footer(s, 10)


# ════════════════════════════════════════════════════════════════════
# S11 — Failure modes (hero figure)
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part II · TrSG")
slide_title(s, "Four Failure Modes — and One Fix",
            "Both window width AND gradient magnitude must scale correctly with V_thr")

# Hero figure: second_main centered.
# Native aspect 0.87 (1905x2180).
# We need final y < 6.30 (badges) so use h = 4.0  ->  w ≈ 3.49.
img_h = 3.7
img_w = img_h * (1905 / 2180)   # ≈ 3.23
img_x = (13.333 - img_w) / 2     # auto-center horizontally
img(s, ASSETS + "/second_main.png", img_x, 2.0, h=img_h)

# Caption directly under the figure
text(s, 0.5, 2.0 + img_h + 0.02, 12.3, 0.25,
     "Each colored box = active gradient region.   width = window  ·  height = magnitude",
     font=SANS, size=9, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Five summary badges in a single horizontal row, BELOW the figure
def badge(slide, x, y, w, h, label, sub, fill, border, sub_color):
    card(slide, x, y, w, h, fill=fill, border=border, border_w=1.5)
    text(slide, x + 0.1, y + 0.1, w - 0.2, 0.35, label,
         font=SANS, size=11, bold=True, color=border, align=PP_ALIGN.CENTER)
    text(slide, x + 0.1, y + 0.45, w - 0.2, 0.4, sub,
         font=SANS, size=9, color=sub_color, align=PP_ALIGN.CENTER)

by  = 6.15
bw  = 2.36
gap = 0.18
bx0 = (13.333 - (5 * bw + 4 * gap)) / 2   # auto-center the badge row

badge(s, bx0 + 0*(bw+gap), by, bw, 0.85,
      "✗ Flood",     "AS-SG · V_thr ≪ 1",
      fill=CORAL_LT, border=CORAL, sub_color=INK)
badge(s, bx0 + 1*(bw+gap), by, bw, 0.85,
      "✗ Starvation", "AS-SG · V_thr ≫ 1",
      fill=CORAL_LT, border=CORAL, sub_color=INK)
badge(s, bx0 + 2*(bw+gap), by, bw, 0.85,
      "✗ Explosion",  "RS-SG · V_thr ≪ 1",
      fill=AMBER_LT, border=AMBER, sub_color=INK)
badge(s, bx0 + 3*(bw+gap), by, bw, 0.85,
      "✗ Vanishing",  "RS-SG · V_thr ≫ 1",
      fill=AMBER_LT, border=AMBER, sub_color=INK)
badge(s, bx0 + 4*(bw+gap), by, bw, 0.85,
      "✓ Balanced",   "TrSG · all regimes",
      fill=TEAL_LT, border=TEAL, sub_color=NAVY)

page_footer(s, 11)


# ════════════════════════════════════════════════════════════════════
# S12 — One-line fix
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part II · TrSG")
slide_title(s, "TrSG — One-Line Fix",
            "Multiply the threshold on the forward path; chain rule cancels the bad factor")

# Top — three step recipe
text(s, 0.5, 2.4, 12.3, 0.4,
     "The recipe",
     font=SANS, size=13, bold=True, color=NAVY)

# Step boxes side by side
y0 = 2.9
def step(slide, x, y, w, h, num, head, eq):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = CREAM
    box.line.color.rgb = TEAL; box.line.width = Pt(1.0)
    text(slide, x + 0.2, y + 0.1, 0.4, 0.4, num,
         font=SANS, size=18, bold=True, color=AMBER)
    text(slide, x + 0.65, y + 0.15, w - 0.7, 0.3, head,
         font=SANS, size=11, bold=True, color=NAVY)
    text(slide, x + 0.65, y + 0.5, w - 0.7, 0.5, eq,
         font="Cambria Math", size=14, color=INK)

step(s, 0.5,  y0, 4.05, 1.05, "1",
     "Relative-scale window", "x = M / V_thr − 1")
step(s, 4.65, y0, 4.05, 1.05, "2",
     "Multiply on forward", "O[t] = V_thr · S[t]")
step(s, 8.8, y0, 4.05, 1.05, "3",
     "Chain rule cancels", "∂O/∂M = f'(x)")

# Chain rule visual
img(s, GEN + "/chain_rule.png", 1.5, 4.3, w=10.3)

# Bottom emphasis
callout_box(s, 0.5, 6.4, 12.3, 0.6,
            fill=TEAL, accent=AMBER,
            title="Window adapts with V_thr  ·  gradient magnitude is threshold-invariant.",
            title_color=WHITE, title_size=13)

page_footer(s, 12)


# ════════════════════════════════════════════════════════════════════
# S13 — Inference stays binary
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part II · TrSG")
slide_title(s, "Inference Stays Binary",
            "TrSG is a training-only trick — neuromorphic hardware never sees it")

# Training row
text(s, 0.5, 2.5, 12.3, 0.4,
     "Training",
     font=SANS, size=14, bold=True, color=NAVY)
callout_box(s, 0.5, 3.0, 12.3, 0.7,
            fill=CREAM, accent=AMBER,
            body="Forward:    O[t]  =  V_thr · S[t]    (real-valued — gradient flows cleanly through V_thr)",
            body_color=INK, body_size=12)

# Inference rows
text(s, 0.5, 3.95, 12.3, 0.4,
     "At deployment, absorb V_thr into the next layer's weights",
     font=SANS, size=14, bold=True, color=NAVY)
callout_box(s, 0.5, 4.45, 12.3, 0.7,
            fill=CREAM, accent=AMBER,
            body="W'_{l→l+1}   =   V_thr^l  ·  W_{l→l+1}    (one-time rescaling per layer)",
            body_color=INK, body_size=12)
callout_box(s, 0.5, 5.25, 12.3, 0.7,
            fill=TEAL_LT, accent=TEAL,
            body="Forward at test time:    O[t]  =  S[t]  ∈  {0, 1}      (standard LIF — back to binary spikes)",
            body_color=NAVY, body_size=12)

# Bottom emphasis
callout_box(s, 0.5, 6.2, 12.3, 0.7,
            fill=TEAL, accent=AMBER,
            title="Standard LIF operators at inference. Hardware compatibility preserved.",
            title_color=WHITE, title_size=13)

page_footer(s, 13)


# ════════════════════════════════════════════════════════════════════
# S14 — TrSG: Result + Mechanism
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part II · TrSG")
slide_title(s, "Why TrSG Works — A Stable Gradient",
            "Three stability metrics tell the same story across extreme V_thr")

# Top: metric definitions
def metric_def(slide, x, y, w, name, formula, healthy):
    card(slide, x, y, w, 1.05, fill=CREAM, border=GRAY_LT, border_w=0.5)
    text(slide, x + 0.15, y + 0.08, w - 0.3, 0.3, name,
         font=SANS, size=12, bold=True, color=NAVY)
    text(slide, x + 0.15, y + 0.4, w - 0.3, 0.3, formula,
         font="Cambria Math", size=11, color=INK)
    text(slide, x + 0.15, y + 0.72, w - 0.3, 0.3, healthy,
         font=SANS, size=9, italic=True, color=GRAY)

metric_def(s, 0.5,  2.45, 4.05,
           "AbsStr",  "mean |gradient|",
           "low → vanishing,   high → exploding")
metric_def(s, 4.65, 2.45, 4.05,
           "RatioAG", "% neurons with active gradient",
           "low → starvation,   high → flooding")
metric_def(s, 8.8,  2.45, 4.05,
           "GradCV",  "std / mean of |gradient|",
           "high → unstable, hard-to-tune updates")

# Middle: stress-test tables (V_thr=0.1, V_thr=2.0)
text(s, 0.5, 3.7, 12.3, 0.3,
     "Stress-test metrics at epoch 100 — CIFAR-100, ResNet-19, τ=2.0 frozen",
     font=SANS, size=10, italic=True, color=GRAY)

def stress_table(slide, x, y, title, rows):
    """Stress-test table using PowerPoint's native table object."""
    # Title bar above the table
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(6.05), Inches(0.42))
    title_bar.fill.solid(); title_bar.fill.fore_color.rgb = NAVY
    title_bar.line.fill.background()
    text(slide, x + 0.2, y + 0.06, 5.8, 0.35, title,
         font=SANS, size=12, bold=True, color=AMBER)

    # Native PPT table:  4 rows (1 header + 3 data) × 4 cols
    tw, th = 6.05, 1.95
    tbl_shape = slide.shapes.add_table(
        4, 4,
        Inches(x), Inches(y + 0.42),
        Inches(tw), Inches(th))
    tbl = tbl_shape.table

    # Column widths (sum to 6.05 inches)
    col_widths = [1.45, 1.50, 1.55, 1.55]
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = Inches(cw)

    # Row heights
    tbl.rows[0].height = Inches(0.45)
    for r in range(1, 4):
        tbl.rows[r].height = Inches(0.50)

    # Header row
    headers = ["SG", "AbsStr", "RatioAG", "GradCV"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = TEAL_LT
        p = cell.text_frame.paragraphs[0]
        p.font.name = SANS; p.font.size = Pt(11); p.font.bold = True
        p.font.color.rgb = NAVY
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)

    # Data rows
    for i, (sg, absstr, ratio, cv, color) in enumerate(rows, start=1):
        # Zebra striping
        bg = CREAM if i % 2 == 0 else WHITE
        for j, val in enumerate([sg, absstr, ratio, cv]):
            cell = tbl.cell(i, j)
            cell.text = val
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            if j == 0:
                p.font.name = SANS; p.font.bold = True
            else:
                p.font.name = "Consolas"; p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = color
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.12)

stress_table(s, 0.5, 4.05,
             "V_thr = 0.1   (small threshold)",
             [("AS-SG",  "0.0145",  "90.20%", "1.22",   CORAL),
              ("RS-SG",  "NaN",     "0.00%",  "NaN",    CORAL),
              ("TrSG",   "0.0153",  "21.45%", "0.87",   TEAL)])

stress_table(s, 6.78, 4.05,
             "V_thr = 2.0   (large threshold)",
             [("AS-SG",  "0.0000",  "0.00%",  "0.000",  CORAL),
              ("RS-SG",  "0.0021",  "0.10%",  "2.25",   CORAL),
              ("TrSG",   "0.0323",  "6.90%",  "0.73",   TEAL)])

# Bottom takeaway
callout_box(s, 0.5, 6.7, 12.3, 0.4,
            fill=NAVY, accent=AMBER,
            title="TrSG keeps all three metrics in the healthy range — at every threshold.",
            title_color=AMBER, title_size=12)

page_footer(s, 14)


# ════════════════════════════════════════════════════════════════════
# S15 — TrSG Robust Across All V_thr  (outcome + real-world relevance)
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Part II · TrSG")
slide_title(s, "TrSG Robust Across All V_thr",
            "Direct accuracy evidence  +  real trainings actually visit the extreme regimes")

# Left: no_training_acc bar chart (the direct accuracy evidence)
text(s, 0.5, 2.4, 6.7, 0.3,
     "Accuracy with V_thr frozen at varying values",
     font=SANS, size=12, bold=True, color=NAVY)
img(s, GEN + "/no_training_acc.png", 0.5, 2.8, w=6.7)
text(s, 0.5, 5.5, 6.7, 0.3,
     "CIFAR-100, ResNet-19, τ=2.0 frozen",
     font=SANS, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Right: threshold trajectories (small) + bullets
text(s, 7.4, 2.4, 5.5, 0.3,
     "Trainable V_thr trajectories — ImageNet",
     font=SANS, size=12, bold=True, color=NAVY)
img(s, ASSETS + "/threshold_training_imagenet.png", 7.4, 2.8, w=5.5, h=2.5)
text(s, 7.4, 5.4, 5.5, 0.3,
     "ResNet-34 on ImageNet · 120 epochs",
     font=SANS, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

# Bullets across the bottom
bullets(s, 0.5, 5.85, 12.3, 0.7, [
    "Frozen V_thr accuracy:  TrSG wins at every threshold; AS-SG/RS-SG fail outside V_thr ≈ 1",
    "Real trainings drift V_thr ≪ 1 in early layers and ≫ 1 in deep layers — exactly the failure regimes",
], size=11, color=INK, bullet="▸")

# Bottom takeaway
callout_box(s, 0.5, 6.65, 12.3, 0.4,
            fill=NAVY, accent=AMBER,
            title="Threshold-robustness is a practical necessity, not a theoretical edge case.",
            title_color=AMBER, title_size=11)

page_footer(s, 15)


# ════════════════════════════════════════════════════════════════════
# S16 — Main Results
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Experiments · 1/3")
slide_title(s, "Main Results — SOTA Across the Board",
            "Standard LIF neurons · no AutoAugment / Cutout · three trials")

# Bar chart on the left
img(s, GEN + "/main_results_bars.png", 0.4, 2.4, w=8.5)

# Three gain cards on right — vertical column
def gain_card(slide, x, y, value, label, color):
    card(slide, x, y, 4.0, 1.4, fill=CREAM, border=color, border_w=1.8)
    text(slide, x + 0.2, y + 0.2, 2.4, 0.6, value,
         font=SERIF, size=26, bold=True, color=color)
    text(slide, x + 2.4, y + 0.32, 1.2, 0.4, "%pt",
         font=SANS, size=11, color=color)
    text(slide, x + 0.2, y + 0.85, 3.7, 0.5, label,
         font=SANS, size=10, color=INK)

gain_card(s, 9.0, 2.4, "+1.09", "CIFAR-100  ·  beats TAB", TEAL)
gain_card(s, 9.0, 3.95, "+0.64", "ImageNet  ·  beats MPS",   TEAL)
gain_card(s, 9.0, 5.5,  "+5.23", "DVS-CIFAR10  ·  beats RMP-Loss", AMBER)

# Bottom emphasis
callout_box(s, 0.5, 6.7, 12.3, 0.4,
            fill=NAVY, accent=AMBER,
            title="Biggest gain on the most challenging (event-based) dataset — DVS-CIFAR10 +5.23 %pt",
            title_color=AMBER, title_size=11)

page_footer(s, 16)


# ════════════════════════════════════════════════════════════════════
# S17 — Ablation
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Experiments · 2/3")
slide_title(s, "Ablation — Independent and Complementary",
            "Each method helps alone; gains stack additively (largest combined effect on event-based data)")

# Bar chart
img(s, GEN + "/ablation_bars.png", 0.4, 2.45, w=10.0)

# Right: takeaways
text(s, 10.8, 2.5, 2.4, 0.4, "Takeaway",
     font=SANS, size=12, bold=True, color=NAVY)

# Three pill takeaways
def pill(slide, x, y, w, h, lab, val, color):
    card(slide, x, y, w, h, fill=CREAM, border=color, border_w=1.3)
    text(slide, x + 0.15, y + 0.1, w - 0.3, 0.3, lab,
         font=SANS, size=10, color=GRAY)
    text(slide, x + 0.15, y + 0.4, w - 0.3, 0.5, val,
         font=SERIF, size=18, bold=True, color=color)

pill(s, 10.6, 2.95, 2.5, 1.0, "MP-Init alone",  "+0.5 ~ 0.8", TEAL)
pill(s, 10.6, 4.05, 2.5, 1.0, "TrSG alone",     "+1.6 ~ 4.2", AMBER)
pill(s, 10.6, 5.15, 2.5, 1.0, "Both",           "+2.1 ~ 4.8", RGBColor(0x0E, 0x5E, 0x5E))

# Bottom message
callout_box(s, 0.5, 6.55, 12.3, 0.55,
            fill=TEAL, accent=AMBER,
            title="The two fixes are orthogonal — diagnoses aim at different components, not the same problem twice.",
            title_color=WHITE, title_size=11)

page_footer(s, 17)


# ════════════════════════════════════════════════════════════════════
# S18 — Generalization
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Experiments · 3/3")
slide_title(s, "Beyond Classification",
            "Transformer SNNs · object detection · energy preserved")

# Three cards
def gen_card(slide, x, y, w, h, header, body):
    card(slide, x, y, w, h, fill=CREAM, border=TEAL, border_w=1.2)
    text(slide, x + 0.2, y + 0.15, w - 0.4, 0.4,
         header, font=SANS, size=12, bold=True, color=TEAL)
    text(slide, x + 0.2, y + 0.55, w - 0.4, h - 0.7, body,
         font=SANS, size=11, color=INK)

gen_card(s, 0.5, 2.5, 6.0, 1.7,
         "Transformer SNN  (QKFormer, ImageNet, T=4)",
         "Baseline 78.80    →    + ours 79.90    (+1.10 %pt)\n"
         "AS-SG    77.09    →    RS-SG  78.70    →    TrSG  79.90\n"
         "Drop-in SG swap on a state-of-the-art Transformer SNN.")

gen_card(s, 6.85, 2.5, 6.0, 1.7,
         "Object detection  (COCO-2017, EMS-ResNet10, T=3)",
         "mAP@0.5         0.291  →  0.305    (+0.014)\n"
         "mAP@0.5:0.95    0.138  →  0.147    (+0.009)\n"
         "First evidence the principles transfer beyond classification.")

gen_card(s, 0.5, 4.4, 6.0, 1.7,
         "Energy under matched configuration  (CIFAR-100)",
         "Baseline   1.57 mJ   ·   75.58 %\n"
         "+ MP-Init + TrSG    1.67 mJ   ·   77.68 %\n"
         "+2.10 %pt at +6 % energy → orders of magnitude better than\n"
         "spike-rate-multiplying baselines (cf. Pareto frontier on slide 9).")

gen_card(s, 6.85, 4.4, 6.0, 1.7,
         "All four surrogate shapes",
         "Rectangular · Triangular · Arctan · Sigmoid\n"
         "TrSG remains the best across initial V_thr, initial τ,\n"
         "and surrogate shape — never tied for second.")

# Bottom emphasis
callout_box(s, 0.5, 6.25, 12.3, 0.75,
            fill=NAVY, accent=AMBER,
            title="No architectural changes, no extra losses, no per-timestep parameters — just drop in.",
            title_color=AMBER, title_size=13)

page_footer(s, 18)


# ════════════════════════════════════════════════════════════════════
# S19 — Recap
# ════════════════════════════════════════════════════════════════════
s = add_blank()
section_label(s, "Conclusion")
slide_title(s, "What This Talk Was About")

# Top: one-line story
callout_box(s, 0.5, 2.4, 12.3, 1.0,
            fill=NAVY, accent=AMBER,
            title="Two long-standing instabilities of direct SNN training, traced and fixed.",
            body="Both fixes are minimal, principled, and preserve standard inference.",
            title_color=AMBER, body_color=WHITE,
            title_size=15, body_size=11)

# Three cards: MP-Init, TrSG, Validated
def takeaway_card(slide, x, y, w, h, header, header_color, body):
    card(slide, x, y, w, h, fill=CREAM, border=header_color, border_w=1.5)
    text(slide, x + 0.2, y + 0.15, w - 0.4, 0.5,
         header, font=SERIF, size=18, bold=True, color=header_color)
    text(slide, x + 0.2, y + 0.7, w - 0.4, h - 0.8,
         body, font=SANS, size=11, color=INK)

takeaway_card(s, 0.5, 3.7, 4.0, 3.0,
              "MP-Init", TEAL,
              "TCS originates in the membrane potential — fix the initial state.\n\n"
              "Per-layer running mean of U[T].\n\n"
              "Zero overhead, standard LIF inference.")

takeaway_card(s, 4.7, 3.7, 4.0, 3.0,
              "TrSG", AMBER,
              "Threshold-induced gradient pathology — fix it with one forward-path multiplication.\n\n"
              "Window adapts, magnitude is invariant.\n\n"
              "Standard LIF after training.")

takeaway_card(s, 8.9, 3.7, 4.0, 3.0,
              "Validated", TEAL,
              "SOTA on CIFAR-10/100, ImageNet, DVS-CIFAR10.\n\n"
              "Generalizes to Transformer SNNs (QKFormer +1.10) and detection (COCO +0.014 mAP).\n\n"
              "All without architectural changes.")

page_footer(s, 19)


# ════════════════════════════════════════════════════════════════════
# S20 — Limitations + Future + Thanks
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, NAVY)
add_circle(s, 9.0, -1.5, 4.0, NAVY_DK)
add_circle(s, -2.0, 5.0, 5.5, NAVY_DK)

text(s, 0.7, 0.7, 10, 0.4, "CONCLUSION",
     font=SANS, size=12, bold=True, color=AMBER, spacing=400)

# Big closing message — different from S18's recap card
text(s, 0.7, 1.6, 12, 0.9,
     "Two instabilities, two principled fixes,",
     font=SERIF, size=36, bold=True, color=WHITE)
text(s, 0.7, 2.5, 12, 0.9,
     "one cohesive story.",
     font=SERIF, size=36, bold=True, color=AMBER)

# Limitations + future — main content of this slide (not duplicate cards)
text(s, 0.7, 4.0, 11, 0.3, "LIMITATIONS",
     font=SANS, size=11, bold=True, color=AMBER, spacing=400)
text(s, 0.7, 4.35, 12.0, 1.1,
     "▸  Per-layer constant approximates π only at the mean — richer parameterizations are open.\n"
     "▸  Hardware with fixed neuron parameters needs quantization-aware adaptation.\n"
     "▸  i.i.d. assumption empirically supported, but tighter bounds welcome.",
     font=SANS, size=12, color=RGBColor(0xCB, 0xD8, 0xE0))

text(s, 0.7, 5.55, 11, 0.3, "FUTURE  WORK",
     font=SANS, size=11, bold=True, color=AMBER, spacing=400)
text(s, 0.7, 5.9, 12.0, 0.9,
     "▸  Hardware-aware adaptation:  Loihi, Akida, TrueNorth\n"
     "▸  Beyond LIF:  adaptive LIF, Izhikevich, foundation-scale SNNs",
     font=SANS, size=12, color=RGBColor(0xCB, 0xD8, 0xE0))

# Amber underline + Thank you (positioned to fit above 7.5" slide bottom)
amber_underline(s, 0.7, 6.85, w=1.0)
text(s, 0.7, 6.97, 12, 0.4,
     "Thank you.   Questions?",
     font=SERIF, size=22, bold=True, italic=True, color=AMBER)
text(s, 11.0, 7.10, 1.8, 0.3, "20 / 20",
     font=SANS, size=10, color=RGBColor(0x80, 0x95, 0xA8), align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════
# APPENDIX — backup slides for Q&A
# ════════════════════════════════════════════════════════════════════
def appendix_footer(slide, label):
    text(slide, 0.5, 7.10, 6, 0.3, "Appendix — Hyunho Kook",
         font=SANS, size=10, color=GRAY)
    text(slide, 7.0, 7.10, 5.8, 0.3, label,
         font=SANS, size=10, color=GRAY, align=PP_ALIGN.RIGHT)

def appendix_label(slide, text_str):
    """Maroon-colored label band in top-left to mark appendix slides."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.62), Inches(0.06), Inches(0.22))
    bar.fill.solid(); bar.fill.fore_color.rgb = CORAL; bar.line.fill.background()
    text(slide, 0.65, 0.6, 6, 0.3, text_str,
         font=SANS, size=10, bold=True, color=CORAL, spacing=400)


# ════════════════════════════════════════════════════════════════════
# A0 — Appendix divider
# ════════════════════════════════════════════════════════════════════
s = add_blank()
fill_bg(s, NAVY)
add_circle(s, -2.0, -2.0, 5.5, NAVY_DK)
add_circle(s, 10.5, 4.5, 4.0, NAVY_DK)

text(s, 0.7, 2.0, 12, 0.4, "BACKUP MATERIAL",
     font=SANS, size=12, bold=True, color=AMBER, spacing=400)
text(s, 0.7, 2.7, 12, 1.2,
     "Appendix",
     font=SERIF, size=60, bold=True, color=WHITE)
amber_underline(s, 0.7, 4.05, w=1.0)
text(s, 0.7, 4.3, 12, 0.5,
     "Q&A reference: hardware, datasets, prior work",
     font=SERIF, size=20, italic=True, color=AMBER)

# Mini table of contents
toc = [
    ("A1", "Can SNNs run on real chips?"),
    ("A2", "Where SNN efficiency comes from"),
    ("A3", "Energy: SNN chip vs DNN platforms"),
    ("A4", "What is an event-based dataset?"),
    ("A5", "From DVS events to SNN input"),
    ("A6", "How were SNNs trained before?"),
    ("A7", "Prior work on TCS & threshold learning"),
]
ty = 5.0
for i, (tag, title) in enumerate(toc):
    col = i // 4
    row = i % 4
    tx = 0.7 + col * 6.3
    yy = ty + row * 0.40
    text(s, tx, yy, 0.6, 0.32, tag,
         font=SANS, size=12, bold=True, color=AMBER)
    text(s, tx + 0.7, yy, 5.6, 0.32, title,
         font=SANS, size=12, color=RGBColor(0xCB, 0xD8, 0xE0))


# ════════════════════════════════════════════════════════════════════
# A1 — SNN chips landscape
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A1")
slide_title(s, "Can SNNs Actually Run on Chips?",
            sub="Production-grade neuromorphic silicon already exists.")

# Left column: brief intro
text(s, 0.5, 2.4, 4.1, 0.4,
     "Yes — and at scale.",
     font=SANS, size=15, bold=True, color=NAVY)
text(s, 0.5, 2.85, 4.1, 2.3,
     "Several research and commercial chips natively run "
     "spiking models. Our methods (MP-Init, TrSG) are "
     "training-time techniques that produce the same LIF model "
     "these chips already accept — so the trained network ports "
     "directly without architectural changes.",
     font=SANS, size=11.5, color=INK)

callout_box(s, 0.5, 5.6, 4.1, 1.3,
            fill=TEAL_LT, accent=TEAL,
            title="Key takeaway",
            body="Hardware is here. The bottleneck is training quality at low T.",
            title_color=TEAL, body_color=INK,
            title_size=11, body_size=11)

# Right column: native PPT table of chips
tx_, ty_, tw_, th_ = 4.95, 2.4, 8.0, 4.5
tbl_shape = s.shapes.add_table(
    8, 5,
    Inches(tx_), Inches(ty_),
    Inches(tw_), Inches(th_))
tbl = tbl_shape.table

col_widths = [1.55, 0.60, 0.85, 1.55, 1.55, 1.90]  # 6 widths for 5 cols (6th will not be set)
col_widths = [1.55, 0.60, 0.95, 1.55, 1.55]  # Sum = 6.20  — actually 5 columns
# Recompute to sum 8.0
col_widths = [1.95, 0.70, 1.05, 1.95, 2.35]
for i, cw in enumerate(col_widths):
    tbl.columns[i].width = Inches(cw)

# Heights
tbl.rows[0].height = Inches(0.50)
for r in range(1, 8):
    tbl.rows[r].height = Inches(0.55)

headers = ["Chip", "Year", "Process", "Capacity", "Headline metric"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.name = SANS; p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.10); cell.margin_right = Inches(0.10)

chips = [
    ("IBM TrueNorth",   "2014", "65 nm",  "1 M neur.\n256 M syn.",  "70 mW · 26 pJ/event"),
    ("Intel Loihi",     "2018", "14 nm",  "130 K neur.\n130 M syn.","23.6 pJ / op"),
    ("Intel Loihi 2",   "2021", "Intel 4","1 M neur.\ngraded spikes","improved · multi-bit"),
    ("SpiNNaker 2",     "2024", "22 nm",  "152 K neur.\n152 M syn. / chip","ARM-based · digital"),
    ("Tianjic",         "2019", "28 nm",  "40 K neur.\nSNN+ANN hybrid","Tsinghua · Nature 2019"),
    ("Speck (SynSense)","2022", "65 nm",  "328 K neur.\nDVS on-chip","< 10 mW · always-on"),
    ("BrainScaleS-2",   "2022", "65 nm",  "512 analog neur.","1000× biological speed"),
]
for i, row in enumerate(chips, start=1):
    bg = CREAM if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        if j == 0:
            p.font.name = SANS; p.font.bold = True; p.font.size = Pt(10.5)
            p.font.color.rgb = NAVY
        elif j == 1 or j == 2:
            p.font.name = "Consolas"; p.font.size = Pt(10)
            p.font.color.rgb = INK
        else:
            p.font.name = SANS; p.font.size = Pt(9.5)
            p.font.color.rgb = INK
        p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        # Multi-line cells
        for extra_p in cell.text_frame.paragraphs[1:]:
            extra_p.font.name = p.font.name
            extra_p.font.size = p.font.size
            extra_p.font.color.rgb = p.font.color.rgb
            extra_p.alignment = p.alignment
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.10); cell.margin_right = Inches(0.10)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)

# Caption under table
text(s, 4.95, 6.95, 8.0, 0.25,
     "Refs: Merolla 2014 · Davies 2018 · Mayr 2019 · Pei 2019 · SynSense 2022 · BrainScaleS team 2022.",
     font=SANS, size=8.5, italic=True, color=GRAY)

appendix_footer(s, "A1 / A7")


# ════════════════════════════════════════════════════════════════════
# A2 — Where SNN efficiency comes from
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A2")
slide_title(s, "Where SNN Efficiency Comes From",
            sub="Three structural advantages — none requires a new training scheme.")

# Three cards across top
card_y = 2.4
card_h = 2.55
card_w = 4.05
gaps = 0.21
xs = [0.5, 0.5 + card_w + gaps, 0.5 + 2*(card_w + gaps)]

reasons = [
    ("①  Event-driven", TEAL,
     "No spike → no compute. Inactive neurons skip MAC and memory access entirely.",
     "Typical activity:  5–20 % per timestep"),
    ("②  Multiplication-free", AMBER,
     "Spike ∈ {0, 1} → W·spike degenerates to W or 0. Each synapse becomes an Accumulate (AC), not a MAC.",
     "INT8 AC ≈ 8× cheaper than INT8 MAC (45 nm)"),
    ("③  Co-located memory", CORAL,
     "Neuromorphic cores keep synaptic weights in local SRAM, on-chip mesh routing — no DRAM trip.",
     "DRAM 64-bit ≈ 1300 pJ vs SRAM ≈ 5 pJ"),
]
for i, (title, accent, body, footer_t) in enumerate(reasons):
    card(s, xs[i], card_y, card_w, card_h, fill=WHITE, border=accent, border_w=1.8)
    # Top accent bar
    abar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(xs[i]), Inches(card_y), Inches(card_w), Inches(0.10))
    abar.fill.solid(); abar.fill.fore_color.rgb = accent; abar.line.fill.background()

    text(s, xs[i] + 0.2, card_y + 0.25, card_w - 0.4, 0.40, title,
         font=SANS, size=14, bold=True, color=accent)
    text(s, xs[i] + 0.2, card_y + 0.75, card_w - 0.4, 1.30, body,
         font=SANS, size=11, color=INK)
    # Bottom strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(xs[i] + 0.15), Inches(card_y + card_h - 0.55),
        Inches(card_w - 0.30), Inches(0.40))
    strip.fill.solid(); strip.fill.fore_color.rgb = CREAM; strip.line.fill.background()
    text(s, xs[i] + 0.2, card_y + card_h - 0.49, card_w - 0.4, 0.32, footer_t,
         font="Consolas", size=10, bold=True, color=NAVY)

# Bottom: AC vs MAC chart  (8.0 × 2.2 figure → at w=7.4 height ≈ 2.04, fits 5.15–7.05)
img(s, f"{GEN}/mac_vs_ac.png", 0.5, 5.15, w=7.4)

# Right side: key fact panel
callout_box(s, 8.20, 5.30, 4.75, 1.65,
            fill=NAVY, accent=AMBER,
            title="Bottom line",
            body="Efficiency is structural — comes from sparsity + binary spikes + co-located memory.\nMP-Init / TrSG protect (1) by stabilizing low-T training.",
            title_color=AMBER, body_color=WHITE,
            title_size=11, body_size=10)

appendix_footer(s, "A2 / A7")


# ════════════════════════════════════════════════════════════════════
# A3 — Energy SNN vs DNN
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A3")
slide_title(s, "Energy: SNN Chip vs DNN Platforms",
            sub="Same task, four hardware platforms — what does the gap actually look like?")

# Left: chart
img(s, f"{GEN}/energy_compare.png", 0.4, 2.4, w=7.7)
text(s, 0.4, 5.95, 7.7, 0.3,
     "Same Aloha keyword-spotting task on each platform.",
     font=SANS, size=10, italic=True, color=GRAY)

# Right: explanation + caveats
text(s, 8.4, 2.4, 4.6, 0.4,
     "How big is the gap?",
     font=SANS, size=15, bold=True, color=NAVY)

text(s, 8.4, 2.85, 4.6, 0.4,
     "Sparse / event tasks", font=SANS, size=12, bold=True, color=TEAL)
text(s, 8.4, 3.18, 4.6, 0.5,
     "~50–100×  vs embedded GPU\n(keyword spot, gesture, sparse coding)",
     font=SANS, size=10.5, color=INK)

text(s, 8.4, 3.95, 4.6, 0.4,
     "Static images (CIFAR / ImageNet)",
     font=SANS, size=12, bold=True, color=AMBER)
text(s, 8.4, 4.28, 4.6, 0.5,
     "~2–5×  vs GPU at iso-accuracy\n(per published SNN/CNN comparisons)",
     font=SANS, size=10.5, color=INK)

text(s, 8.4, 4.85, 4.6, 0.4,
     "Why so task-dependent?",
     font=SANS, size=12, bold=True, color=CORAL)
text(s, 8.4, 5.18, 4.6, 0.7,
     "Energy ∝ T × spike-rate × neurons.\nLow-T + high-sparsity → big gap.",
     font=SANS, size=10.5, color=INK)

callout_box(s, 8.4, 6.05, 4.6, 0.9,
            fill=NAVY, accent=AMBER,
            title="Why this matters for our work",
            body="MP-Init enables T = 2 SOTA → directly multiplies energy savings.",
            title_color=AMBER, body_color=WHITE,
            title_size=10.5, body_size=9.5)

# Citation strip
text(s, 0.5, 7.0, 12.0, 0.25,
     "Sources: Blouw et al. 2019 (keyword-spotting benchmark); Davies et al. 2021 (LASSO 48× vs CPU); Horowitz ISSCC 2014 (per-op energy).",
     font=SANS, size=8.5, italic=True, color=GRAY)

appendix_footer(s, "A3 / A7")


# ════════════════════════════════════════════════════════════════════
# A4 — What is an event-based dataset?
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A4")
slide_title(s, "What Is an Event-Based Dataset?",
            sub="Output of asynchronous DVS cameras — sparse, fast, high dynamic range.")

# Left column: DVS explanation
text(s, 0.5, 2.4, 5.5, 0.4,
     "DVS: Dynamic Vision Sensor",
     font=SANS, size=15, bold=True, color=NAVY)
text(s, 0.5, 2.85, 5.5, 1.2,
     "Each pixel asynchronously emits an event "
     "when its log-intensity changes by a threshold.\n"
     "Output: a stream of (x, y, t, polarity) tuples.",
     font=SANS, size=11, color=INK)

# 4 properties as small boxes
prop_y = 4.35
props = [
    ("μs",  "temporal resolution", TEAL),
    (">120 dB", "dynamic range", AMBER),
    ("Sparse", "only changes emit", CORAL),
    ("No blur","no exposure window", NAVY),
]
for i, (big, small, color) in enumerate(props):
    bx = 0.5 + i * 1.42
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(bx), Inches(prop_y), Inches(1.30), Inches(0.95))
    box.fill.solid(); box.fill.fore_color.rgb = CREAM
    box.line.color.rgb = color; box.line.width = Pt(1.5)
    text(s, bx, prop_y + 0.10, 1.30, 0.40, big,
         font=SANS, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    text(s, bx, prop_y + 0.50, 1.30, 0.40, small,
         font=SANS, size=8.5, color=GRAY, align=PP_ALIGN.CENTER)

# This work uses
callout_box(s, 0.5, 5.65, 5.85, 1.10,
            fill=TEAL_LT, accent=TEAL,
            title="In this thesis",
            body="DVS-CIFAR10  (event-based image classification, T = 10).",
            title_color=TEAL, body_color=INK,
            title_size=11, body_size=11)

# Right column: native PPT table of common datasets
tx_, ty_, tw_, th_ = 6.6, 2.4, 6.4, 4.35
tbl_shape = s.shapes.add_table(
    8, 3,
    Inches(tx_), Inches(ty_),
    Inches(tw_), Inches(th_))
tbl = tbl_shape.table

col_widths = [2.30, 2.85, 1.25]
for i, cw in enumerate(col_widths):
    tbl.columns[i].width = Inches(cw)

tbl.rows[0].height = Inches(0.45)
for r in range(1, 8):
    tbl.rows[r].height = Inches(0.55)

headers = ["Dataset", "Description", "Size"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.name = SANS; p.font.size = Pt(11); p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT if j != 2 else PP_ALIGN.CENTER
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)

datasets = [
    ("DVS-CIFAR10",      "CIFAR-10 on monitor + saccading DVS",     "10 K"),
    ("N-MNIST",          "MNIST shown to event camera",              "70 K"),
    ("N-Caltech101",     "Caltech101 in event form",                "9.1 K"),
    ("DVS-Gesture (IBM)","11 hand gestures, 29 subjects",          "1.3 K clips"),
    ("SHD",              "Spiking Heidelberg Digits (audio events)", "10 K"),
    ("Gen1 / 1Mpx",      "Prophesee automotive (object detect.)",   "39 h / 14 h"),
    ("ASL-DVS",          "American Sign Language (event)",          "100 K"),
]
for i, (name, desc, sz) in enumerate(datasets, start=1):
    bg = CREAM if i % 2 == 0 else WHITE
    is_used = (name == "DVS-CIFAR10")
    for j, val in enumerate([name, desc, sz]):
        cell = tbl.cell(i, j)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = (TEAL_LT if is_used else bg)
        p = cell.text_frame.paragraphs[0]
        if j == 0:
            p.font.name = SANS; p.font.bold = True; p.font.size = Pt(10.5)
            p.font.color.rgb = (TEAL if is_used else NAVY)
        elif j == 1:
            p.font.name = SANS; p.font.size = Pt(10)
            p.font.color.rgb = INK
        else:
            p.font.name = "Consolas"; p.font.size = Pt(10)
            p.font.color.rgb = INK
        p.alignment = PP_ALIGN.LEFT if j != 2 else PP_ALIGN.CENTER
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)

text(s, 6.6, 6.8, 6.4, 0.25,
     "Highlighted: dataset used in this thesis.",
     font=SANS, size=9, italic=True, color=GRAY)

appendix_footer(s, "A4 / A7")


# ════════════════════════════════════════════════════════════════════
# A5 — From DVS events to SNN input
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A5")
slide_title(s, "From DVS Events to SNN Input",
            sub="Asynchronous events become T frames that the LIF stack consumes.")

# Hero diagram (large) — captions are baked into the matplotlib figure itself
img(s, f"{GEN}/dvs_workflow.png", 0.6, 2.5, w=12.1)

# Bottom callout
callout_box(s, 0.6, 6.30, 12.1, 0.65,
            fill=TEAL_LT, accent=TEAL,
            title="Why this matters for MP-Init",
            body="The first frame after binning has no history → the membrane starts at zero, far from π. That's exactly the gap MP-Init closes.",
            title_color=TEAL, body_color=INK,
            title_size=11.5, body_size=10.5)

appendix_footer(s, "A5 / A7")


# ════════════════════════════════════════════════════════════════════
# A6 — How were SNNs trained before?
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A6")
slide_title(s, "How Were SNNs Trained Before?",
            sub="From biological learning rules to modern surrogate-gradient BPTT.")

# Timeline backbone
tl_y = 3.2
tl_x_start = 0.8
tl_x_end = 12.5
# Main horizontal bar
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(tl_x_start), Inches(tl_y), Inches(tl_x_end - tl_x_start), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()

# Era markers
eras = [
    # (x_pos, year, era_name, accent_color)
    (0.8,  "1998", "STDP",                  AMBER),
    (4.0,  "2015", "ANN→SNN\nconversion",   CORAL),
    (8.0,  "2018", "Surrogate\ngradient",   TEAL),
    (12.0, "today","Direct +\nlearnable",   NAVY),
]

for (xp, year, name, color) in eras:
    # Tick + circle on timeline
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
        Inches(xp - 0.13), Inches(tl_y - 0.10), Inches(0.26), Inches(0.26))
    circ.fill.solid(); circ.fill.fore_color.rgb = color
    circ.line.color.rgb = WHITE; circ.line.width = Pt(2.5)

    # Year above
    text(s, xp - 0.6, tl_y - 0.5, 1.2, 0.3, year,
         font="Consolas", size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    # Era name below
    text(s, xp - 1.0, tl_y + 0.25, 2.0, 0.6, name,
         font=SANS, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)

# Detail cards under timeline (one per era)
detail_y = 4.55
detail_h = 2.10
details = [
    ("Spike-Timing-Dependent Plasticity",
     "Biological Hebbian rule.\nUnsupervised, local; works on small nets only.",
     "Bi & Poo 1998\nDiehl & Cook 2015",
     0.55, 2.85, AMBER),
    ("Train ANN → convert weights",
     "Use pretrained ReLU CNN; map activations to spike rates.\nEarly: T = 100–2000.  Modern (QCFS '23): T ≤ 8.\nStill cannot consume native event data.",
     "Cao 2015 · Rueckauer 2017\nBu 2023 (QCFS) · Hao 2023",
     3.55, 2.85, CORAL),
    ("Surrogate gradient (STBP / BPTT)",
     "Replace Heaviside backward with a smooth function.\nFew timesteps; works on event data; modern dominant.",
     "Bohte 2002 · Wu 2018\nNeftci 2019",
     6.55, 2.85, TEAL),
    ("Learnable LIF parameters",
     "Train V_thr, τ end-to-end → adaptive neurons.\nThis thesis: stabilize that training pathway.",
     "Fang 2021 · Wang 2022\nRathi 2023 · ours",
     9.55, 2.85, NAVY),
]
for (title, body, refs, x, w, color) in details:
    card(s, x, detail_y, w, detail_h, fill=WHITE, border=color, border_w=1.5)
    # accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(detail_y), Inches(w), Inches(0.06))
    strip.fill.solid(); strip.fill.fore_color.rgb = color; strip.line.fill.background()

    text(s, x + 0.15, detail_y + 0.18, w - 0.3, 0.4, title,
         font=SANS, size=11, bold=True, color=color)
    text(s, x + 0.15, detail_y + 0.62, w - 0.3, 1.05, body,
         font=SANS, size=10, color=INK)
    # Refs at bottom in italic
    text(s, x + 0.15, detail_y + detail_h - 0.55, w - 0.3, 0.5, refs,
         font=SANS, size=8.5, italic=True, color=GRAY)

appendix_footer(s, "A6 / A7")


# ════════════════════════════════════════════════════════════════════
# A7 — Prior work on TCS & threshold learning
# ════════════════════════════════════════════════════════════════════
s = add_blank()
appendix_label(s, "APPENDIX · A7")
slide_title(s, "Where Prior Work Falls Short",
            sub="TCS methods miss the membrane potential; threshold-learning methods don't diagnose the gradient.")

# ── Left: TCS prior work table ────────────────────────────────────
text(s, 0.5, 2.4, 6.0, 0.4,
     "TCS / Normalization in SNNs",
     font=SANS, size=14, bold=True, color=NAVY)

tx_, ty_, tw_, th_ = 0.5, 2.85, 6.0, 3.0
tbl_shape = s.shapes.add_table(
    5, 3,
    Inches(tx_), Inches(ty_),
    Inches(tw_), Inches(th_))
tbl = tbl_shape.table

col_widths = [1.40, 2.55, 2.05]
for i, cw in enumerate(col_widths):
    tbl.columns[i].width = Inches(cw)
tbl.rows[0].height = Inches(0.50)
for r in range(1, 5):
    tbl.rows[r].height = Inches(0.62)

headers = ["Method", "What it does", "What it misses"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.name = SANS; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)

tcs_rows = [
    ("BNTT '21",   "BN on synaptic input current",  "U[t] drift not corrected"),
    ("tdBN '20",   "Threshold-dependent BN on activations", "Pre-spike BN; U[t] persists"),
    ("TEBN '22",   "Timestep-specific affine on outputs", "Membrane drift untouched"),
    ("TAB '24",    "Temporal-adaptive output normalization",  "Membrane TCS persists"),
]
for i, (m, w_, mi) in enumerate(tcs_rows, start=1):
    bg = CREAM if i % 2 == 0 else WHITE
    for j, val in enumerate([m, w_, mi]):
        cell = tbl.cell(i, j)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        if j == 0:
            p.font.name = SANS; p.font.size = Pt(10); p.font.bold = True
            p.font.color.rgb = NAVY
        elif j == 2:
            p.font.name = SANS; p.font.size = Pt(9.5); p.font.italic = True
            p.font.color.rgb = CORAL
        else:
            p.font.name = SANS; p.font.size = Pt(9.5)
            p.font.color.rgb = INK
        p.alignment = PP_ALIGN.LEFT
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)

callout_box(s, 0.5, 6.0, 6.0, 0.85,
            fill=TEAL_LT, accent=TEAL,
            title="Our take (MP-Init)",
            body="Fix the cause, not symptom — align M[t] with stationary π directly.",
            title_color=TEAL, body_color=INK,
            title_size=11, body_size=10)

# ── Right: Threshold-learning prior work table ────────────────────
text(s, 6.85, 2.4, 6.0, 0.4,
     "Trainable LIF parameters (V_thr, τ)",
     font=SANS, size=14, bold=True, color=NAVY)

tx_, ty_, tw_, th_ = 6.85, 2.85, 6.0, 3.0
tbl_shape = s.shapes.add_table(
    5, 3,
    Inches(tx_), Inches(ty_),
    Inches(tw_), Inches(th_))
tbl = tbl_shape.table

col_widths = [1.55, 2.40, 2.05]
for i, cw in enumerate(col_widths):
    tbl.columns[i].width = Inches(cw)
tbl.rows[0].height = Inches(0.50)
for r in range(1, 5):
    tbl.rows[r].height = Inches(0.62)

headers = ["Method", "What's learned", "What it misses"]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    p = cell.text_frame.paragraphs[0]
    p.font.name = SANS; p.font.size = Pt(10.5); p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)

thr_rows = [
    ("PLIF '21",     "Leakage τ only",            "V_thr fixed"),
    ("LTMD '22",     "V_thr + τ (heavy clipping)", "Hides instability"),
    ("DIET-SNN '23", "V_thr + τ (gradient clip)", "Doesn't classify failure"),
    ("Ours (TrSG)",  "V_thr + τ, scale-invariant", "—"),
]
for i, (m, w_, mi) in enumerate(thr_rows, start=1):
    is_ours = (m == "Ours (TrSG)")
    bg = (TEAL_LT if is_ours else (CREAM if i % 2 == 0 else WHITE))
    for j, val in enumerate([m, w_, mi]):
        cell = tbl.cell(i, j)
        cell.text = val
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        if j == 0:
            p.font.name = SANS; p.font.size = Pt(10); p.font.bold = True
            p.font.color.rgb = (TEAL if is_ours else NAVY)
        elif j == 2:
            if is_ours:
                p.font.name = SANS; p.font.size = Pt(11); p.font.bold = True
                p.font.color.rgb = TEAL
            else:
                p.font.name = SANS; p.font.size = Pt(9.5); p.font.italic = True
                p.font.color.rgb = CORAL
        else:
            p.font.name = SANS; p.font.size = Pt(9.5)
            p.font.color.rgb = (TEAL if is_ours else INK)
            if is_ours: p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.12); cell.margin_right = Inches(0.10)

callout_box(s, 6.85, 6.0, 6.0, 0.85,
            fill=NAVY, accent=AMBER,
            title="Our take (TrSG)",
            body="Diagnose first (4 failure modes), then cancel the 1/V_thr factor by construction.",
            title_color=AMBER, body_color=WHITE,
            title_size=11, body_size=10)

appendix_footer(s, "A7 / A7")


# ─── Save ──────────────────────────────────────────────────────────
out = "Hyunho_Kook_Master_Defense.pptx"
prs.save(out)
print(f"Saved: {out}    Slides: {len(prs.slides)}")
